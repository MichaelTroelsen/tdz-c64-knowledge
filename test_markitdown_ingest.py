"""markitdown-fallback ingestion tests.

kb/ingest/_extraction.py routes an EXPLICIT allowlist of extensions
(.docx/.pptx/.epub/.csv/.json/.xml) to markitdown, as a fallback for formats
no native handler claims. Every native format (.pdf/.xlsx/.zip/...) must stay
on its own extractor - a later refactor that widens the allowlist or drops
the elif's exclusivity would silently route a native format's bytes through
markitdown instead. That is the property the headline test below pins.

markitdown, magika and onnxruntime ARE installed in this venv (see CLAUDE.md),
so features.MARKITDOWN_SUPPORT is True and this suite genuinely exercises the
markitdown path rather than skipping around it.
"""

import io
import json
import os
import zipfile

import pytest

import features
from models import UnsupportedFileTypeError

pytestmark = pytest.mark.skipif(
    not features.MARKITDOWN_SUPPORT,
    reason="markitdown not installed - see CLAUDE.md",
)


@pytest.fixture
def kb(tmp_path, monkeypatch):
    """A KnowledgeBase on an isolated data dir, never the live one."""
    monkeypatch.setenv("TDZ_DATA_DIR", str(tmp_path))
    monkeypatch.setenv("USE_BM25", "0")
    from server import KnowledgeBase
    instance = KnowledgeBase(str(tmp_path))
    yield instance
    instance.close()


# --- fixture builders --------------------------------------------------------

_CONTENT_TYPES_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
    '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
    '<Default Extension="xml" ContentType="application/xml"/>'
    '<Override PartName="/word/document.xml" '
    'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
    '</Types>'
)

_RELS_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
    '<Relationship Id="rId1" '
    'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
    'Target="word/document.xml"/>'
    '</Relationships>'
)

_DOCUMENT_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
    '<w:body><w:p><w:r><w:t>Commando by Rob Hubbard</w:t></w:r></w:p></w:body>'
    '</w:document>'
)


def _write_docx(path):
    """Hand-assemble a minimal OPC docx (mammoth, which markitdown uses, only
    needs these three members: content types, package rels, and the body)."""
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("[Content_Types].xml", _CONTENT_TYPES_XML)
        zf.writestr("_rels/.rels", _RELS_XML)
        zf.writestr("word/document.xml", _DOCUMENT_XML)
    return path


def _write_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Commando"
    slide.placeholders[1].text = "Rob Hubbard"
    prs.save(str(path))
    return path


def _write_xlsx(path):
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws["A1"] = "native-xlsx-sentinel"
    wb.save(str(path))
    return path


def _write_pdf(path):
    """The smallest PDF pypdf will open: one empty page, no content stream."""
    pdf_bytes = (
        b"%PDF-1.4\n"
        b"1 0 obj<</Type/Catalog/Pages 2 0 R>>endobj\n"
        b"2 0 obj<</Type/Pages/Kids[3 0 R]/Count 1>>endobj\n"
        b"3 0 obj<</Type/Page/Parent 2 0 R/MediaBox[0 0 200 200]>>endobj\n"
        b"xref\n0 4\n0000000000 65535 f \n"
        b"trailer<</Size 4/Root 1 0 R>>\n"
        b"startxref\n0\n%%EOF"
    )
    path.write_bytes(pdf_bytes)
    return path


def _write_sid_zip(path):
    with zipfile.ZipFile(path, "w") as zf:
        zf.writestr("readme.txt", "not a tune")
    return path


# --- 1. positive path: real docx / pptx -------------------------------------

def test_docx_ingests_via_markitdown(kb, tmp_path):
    f = _write_docx(tmp_path / "sample.docx")
    text, file_type, total_pages, pdf_metadata = kb._extract_text_for_file(str(f))

    assert file_type == "docx"
    assert total_pages is None
    assert "Commando" in text


def test_pptx_ingests_via_markitdown(kb, tmp_path):
    f = _write_pptx(tmp_path / "sample.pptx")
    text, file_type, total_pages, pdf_metadata = kb._extract_text_for_file(str(f))

    assert file_type == "pptx"
    assert total_pages is None
    assert "Commando" in text
    assert "Rob Hubbard" in text


# --- 2. the library-free allowlist members -----------------------------------

def test_csv_ingests_via_markitdown(kb, tmp_path):
    f = tmp_path / "sample.csv"
    f.write_text("title,author\nCommando,Rob Hubbard\n", encoding="utf-8")
    text, file_type, total_pages, pdf_metadata = kb._extract_text_for_file(str(f))

    assert file_type == "csv"
    assert total_pages is None
    assert "Commando" in text


def test_json_ingests_via_markitdown(kb, tmp_path):
    f = tmp_path / "sample.json"
    f.write_text(json.dumps({"title": "Commando", "author": "Rob Hubbard"}), encoding="utf-8")
    text, file_type, total_pages, pdf_metadata = kb._extract_text_for_file(str(f))

    assert file_type == "json"
    assert total_pages is None
    assert "Commando" in text


def test_xml_ingests_via_markitdown(kb, tmp_path):
    f = tmp_path / "sample.xml"
    f.write_text("<tune><title>Commando</title><author>Rob Hubbard</author></tune>", encoding="utf-8")
    text, file_type, total_pages, pdf_metadata = kb._extract_text_for_file(str(f))

    assert file_type == "xml"
    assert total_pages is None
    assert "Commando" in text


# --- 3. negative: an unknown extension is refused, generically --------------

def test_unknown_extension_raises_generic_unsupported_error(kb, tmp_path):
    f = tmp_path / "sample.qqq"
    f.write_text("whatever", encoding="utf-8")

    with pytest.raises(UnsupportedFileTypeError) as exc:
        kb._extract_text_for_file(str(f))

    message = str(exc.value)
    assert "Unsupported file type: .qqq" in message
    # the fallback must not have swallowed this into the markitdown message
    assert "markitdown" not in message.lower()


# --- 4. disabled path: TDZ_MARKITDOWN=0 --------------------------------------

def test_disabled_markitdown_refuses_docx_naming_the_extra(kb, tmp_path, monkeypatch):
    monkeypatch.setenv("TDZ_MARKITDOWN", "0")
    f = _write_docx(tmp_path / "sample.docx")

    with pytest.raises(UnsupportedFileTypeError) as exc:
        kb._extract_text_for_file(str(f))

    message = str(exc.value)
    assert ".docx" in message
    assert "markitdown" in message.lower()
    assert "TDZ_MARKITDOWN=0" in message


# --- 5. THE HEADLINE ASSERTION: native formats never reach markitdown -------

@pytest.mark.parametrize("build,filename,expected_type", [
    (_write_pdf, "sample.pdf", "pdf"),
    (_write_sid_zip, "sample.zip", None),   # zip has no readable tune - fine, it must still not touch markitdown
    (_write_xlsx, "sample.xlsx", "excel"),
])
def test_native_formats_never_consult_markitdown(kb, tmp_path, monkeypatch, build, filename, expected_type):
    calls = []

    def _spy(self, filepath):
        calls.append(filepath)
        raise AssertionError(
            f"_extract_via_markitdown was called for a native format: {filepath}"
        )

    monkeypatch.setattr(
        "kb.ingest._extraction._ExtractionMixin._extract_via_markitdown", _spy
    )

    f = build(tmp_path / filename)
    try:
        text, file_type, total_pages, pdf_metadata = kb._extract_text_for_file(str(f))
    except Exception:
        # native extraction may legitimately fail on a minimal fixture (e.g.
        # the sid-less zip) - what matters is that markitdown was never
        # reached, not that extraction succeeded.
        pass
    else:
        if expected_type is not None:
            assert file_type == expected_type

    assert calls == [], f"markitdown was consulted for {filename}: {calls}"


# --- 6. the converter is constructed once and cached -------------------------

def test_markitdown_converter_is_constructed_once_and_cached(kb, tmp_path):
    f1 = tmp_path / "one.docx"
    f2 = tmp_path / "two.docx"
    _write_docx(f1)
    _write_docx(f2)

    assert getattr(kb, "_markitdown_converter", None) is None

    kb._extract_text_for_file(str(f1))
    converter_after_first = kb._markitdown_converter
    assert converter_after_first is not None

    kb._extract_text_for_file(str(f2))
    converter_after_second = kb._markitdown_converter

    assert converter_after_second is converter_after_first
